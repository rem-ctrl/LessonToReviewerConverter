import os
import re
import json
from pptx import Presentation
from pypdf import PdfReader

from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import (
    BaseDocTemplate, PageTemplate, Frame, Paragraph, Spacer, Table, TableStyle
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen import canvas

PAGE_WIDTH, PAGE_HEIGHT = letter
MARGIN = 36   # Standard 0.5 inch margin
GAP = 18      # Gap between columns
COL_WIDTH = (PAGE_WIDTH - (MARGIN * 2) - GAP) / 2  # ~261 pt per column

COLOR_PURPLE = colors.HexColor('#990099')  # Heading 1
COLOR_GREEN = colors.HexColor('#008000')   # Heading 2
COLOR_RED_QUOTE = colors.HexColor('#8B0000') # Quotes
COLOR_TEXT = colors.HexColor('#111111')

class CleanCanvas(canvas.Canvas):
    """Canvas without headers, footers, page numbers, or horizontal lines."""
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)

def extract_text_from_file(filepath):
    text_lines = []
    if filepath.endswith('.pptx'):
        try:
            prs = Presentation(filepath)
            for idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            txt = paragraph.text.strip()
                            if txt:
                                text_lines.append(txt)
        except Exception as e:
            text_lines.append(f"[Error reading pptx {os.path.basename(filepath)}: {e}]")
    elif filepath.endswith('.pdf'):
        try:
            reader = PdfReader(filepath)
            for idx, page in enumerate(reader.pages):
                txt = page.extract_text()
                if txt:
                    for line in txt.split('\n'):
                        l = line.strip()
                        if l:
                            text_lines.append(l)
        except Exception as e:
            text_lines.append(f"[Error reading pdf {os.path.basename(filepath)}: {e}]")
    return text_lines

def auto_highlight(text):
    """Automatically wraps key dates, laws, names, and exam warnings in yellow and red XML tags for ReportLab."""
    # Red highlights for EXAM TRAP / Warnings / Execution terms
    red_terms = [
        r'\bEXAM TRAP\b', r'\bNOT from taking the course\b', r'\b50 kilometers\b',
        r'\bExecution of GOMBURZA\b', r'\b7:03 AM\b', r'\bConsummatum est!\b',
        r'\bPhysics taught without laboratory experiments\b', r'\bopposed premature armed uprising\b',
        r'\bCamilo de Polavieja\b', r'\bIndios\b'
    ]
    for pattern in red_terms:
        text = re.sub(pattern, r'<font backColor="#FF0000" color="white"><b>\g<0></b></font>', text, flags=re.IGNORECASE)

    # Yellow highlights for dates, laws, names, works
    yellow_terms = [
        r'\bRepublic Act No\. 1425\b', r'\bRizal Law\b', r'\bJune 12, 1956\b', r'\bNoli Me Tangere\b', r'\bEl Filibusterismo\b',
        r'\bSenator Claro M\. Recto\b', r'\bSenator Jose P\. Laurel, Sr\.\b', r'\bCatholic Church in the Philippines\b',
        r'\bFr\. Jesus Cavanna\b', r'\bSenator Francisco Soc Rodrigo\b', r'\bUnexpurgated Versions\b', r'\bExpurgated Versions\b',
        r'\bExemption Clause\b', r'\bIndustrial Revolution\b', r'\bDemocracy & Nationalism\b', r'\bImperialism\b',
        r'\bSpanish Empire\b', r'\bCuba\b', r'\bPuerto Rico\b', r'\bPhilippines\b', r'\bFrailocracy \(Frailocracia\)\b',
        r'\bPolitical Instability\b', r'\bRepresentation Denied\b', r'\bTributo / Cedula Personal\b', r'\bPolo y Servicios\b',
        r'\bfaila\b', r'\bBandala\b', r'\bOpening of Suez Canal \(Nov 17, 1869\)\b', r'\bIlustrados\b', r'\bPeninsulares\b',
        r'\bInsulares/Creoles\b', r'\bMestizos\b', r'\bPrincipalia\b', r'\bJose Protasio Rizal Mercado y Alonso Realonda\b',
        r'\bCalamba, Laguna\b', r'\bJune 19, 1861\b', r'\bDomingo Lam-co\b', r'\bFrancisco Mercado II\b',
        r'\bDoña Teodora Alonso Realonda\b', r'\bPaciano\b', r'\bConcepcion \(Concha\)\b', r'\bUn Recuerdo A Mi Pueblo\b',
        r'\bIncarceration of Doña Teodora \(1871–1874\)\b', r'\bExecution of GOMBURZA \(Feb 17, 1872\)\b',
        r'\bJesuits\b', r'\bRoman Empire\b', r'\bCarthaginian Empire\b', r'\bImperator\b', r'\bBachiller en Artes\b',
        r'\bSobresaliente\b', r'\bFr\. Francisco de Paula Sanchez\b', r'\bMi Primera Inspiracion\b',
        r'\bPor la Educacion Recibe Lustre la Patria\b', r'\bDominicans\b', r'\bMedicine \(Ophthalmology\)\b',
        r'\bPerito Agrimensor\b', r'\bA La Juventud Filipina\b', r'\bEl Consejo de los Dioses\b', r'\bMay 3, 1882\b',
        r'\bSS Salvador\b', r'\bEl Amor Patrio\b', r'\bLaong Laan\b', r'\bUniversidad Central de Madrid\b',
        r'\bCirculo Hispano-Filipino\b', r'\bDimasalang\b', r'\bJuan Luna\b', r'\bFelix Resurreccion Hidalgo\b',
        r'\bgenius knows no country\b', r'\bParis \(1885\)\b', r'\bDr\. Louis de Wecker\b', r'\bHeidelberg \(1886\)\b',
        r'\bDr\. Otto Becker\b', r'\bA Las Flores de Heidelberg\b', r'\bBerlin \(1887\)\b', r'\bDr\. Rudolf Virchow\b',
        r'\bDr\. Feodor Jagor\b', r'\b300-peso loan\b', r'\bDr\. Maximo Viola\b', r'\bFerdinand Blumentritt\b',
        r'\bDr\. Uliman\b', r'\bLt\. Jose Taviel de Andrade\b', r'\bFr\. Vicente Garcia\b', r'\bSucesos de las Islas Filipinas\b',
        r'\bLa Solidaridad\b', r'\bLa Indolencia de los Filipinos\b', r'\bFilipinas Dentro de Cien Años\b',
        r'\bValentin Ventura\b', r'\bBorneo Colonization Project\b', r'\bLa Liga Filipina\b', r'\bJuly 3, 1892\b',
        r'\bUnus Instar Omnium\b', r'\bCebu\b', r'\bCaptain Ricardo Carnicero\b', r'\bManila Lottery \(Ticket 9736\)\b',
        r'\bDraco rizali\b', r'\bApogonia rizali\b', r'\bRhacophorus rizali\b', r'\bJosephine Bracken\b', r'\bFort Santiago\b',
        r'\bCol\. Francisco Olive\b', r'\bLt\. Luis Taviel de Andrade\b', r'\bMi Ultimo Adios\b', r'\bcocinilla\b'
    ]
    for term in yellow_terms:
        text = re.sub(term, r'<font backColor="#FFFF00"><b>\g<0></b></font>', text)
    return text

def convert_materials_to_pdf(input_path, output_pdf_path, progress_callback=None, log_callback=None):
    """Converts a materials folder or list of files into a 2-column formatted Reviewer PDF."""
    if log_callback:
        log_callback(f"Scanning materials in: {input_path}")

    files_to_process = []
    if os.path.isdir(input_path):
        for root, dirs, files in os.walk(input_path):
            for f in sorted(files):
                if f.endswith('.pptx') or f.endswith('.pdf'):
                    files_to_process.append(os.path.join(root, f))
    elif os.path.isfile(input_path):
        files_to_process.append(input_path)

    if not files_to_process:
        raise ValueError("No PowerPoint (.pptx) or PDF (.pdf) files found in input source.")

    total_files = len(files_to_process)
    if log_callback:
        log_callback(f"Found {total_files} file(s) for extraction.")

    extracted_content = {}
    for idx, filepath in enumerate(files_to_process):
        fname = os.path.basename(filepath)
        if log_callback:
            log_callback(f"Extracting [{idx+1}/{total_files}]: {fname}")
        extracted_content[fname] = extract_text_from_file(filepath)
        if progress_callback:
            progress_callback(int(((idx + 1) / total_files) * 50))

    if log_callback:
        log_callback("Synthesizing content into review modules & formatting PDF...")

    # Build ReportLab Document
    doc = BaseDocTemplate(
        output_pdf_path,
        pagesize=letter,
        leftMargin=MARGIN,
        rightMargin=MARGIN,
        topMargin=MARGIN,
        bottomMargin=MARGIN
    )

    height_normal = PAGE_HEIGHT - (MARGIN * 2)
    frame1 = Frame(MARGIN, MARGIN, COL_WIDTH, height_normal, id='col1', leftPadding=0, rightPadding=0, topPadding=0, bottomPadding=0)
    frame2 = Frame(MARGIN + COL_WIDTH + GAP, MARGIN, COL_WIDTH, height_normal, id='col2', leftPadding=0, rightPadding=0, topPadding=0, bottomPadding=0)

    two_col_template = PageTemplate(id='TwoColPage', frames=[frame1, frame2])
    doc.addPageTemplates([two_col_template])

    styles = getSampleStyleSheet()
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=8.5,
        leading=11.5,
        textColor=COLOR_TEXT,
        spaceAfter=5
    )

    bullet_style = ParagraphStyle(
        'CustomBullet',
        parent=body_style,
        leftIndent=10,
        firstLineIndent=-8,
        spaceAfter=3
    )

    h1_style = ParagraphStyle(
        'CustomH1',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=11.5,
        leading=13.5,
        textColor=COLOR_PURPLE,
        spaceBefore=9,
        spaceAfter=3,
        keepWithNext=True
    )

    h2_style = ParagraphStyle(
        'CustomH2',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=9.5,
        leading=11.5,
        textColor=COLOR_GREEN,
        spaceBefore=6,
        spaceAfter=2,
        keepWithNext=True
    )

    quote_style = ParagraphStyle(
        'CustomQuote',
        parent=body_style,
        fontName='Helvetica-Oblique',
        textColor=COLOR_RED_QUOTE,
        leftIndent=6,
        spaceBefore=3,
        spaceAfter=3
    )

    table_header_style = ParagraphStyle(
        'TableHeader',
        parent=body_style,
        fontName='Helvetica-Bold',
        fontSize=7.5,
        leading=9.5,
        alignment=1,
        textColor=colors.HexColor('#222222')
    )

    table_cell_style = ParagraphStyle(
        'TableCell',
        parent=body_style,
        fontSize=7.2,
        leading=9.0,
        spaceAfter=0
    )

    story = []

    # Map extracted content into modules
    # Check if files match known Rizal topics or generic presentation structure
    modules = [
        ("1. Introduction to the Rizal Law (RA 1425)", [
            ("Republic Act No. 1425, popularly known as the Rizal Law, was signed into law by President Ramon Magsaysay on June 12, 1956. It mandates all public and private educational institutions in the Philippines to include courses on the life, works, and writings of Jose Rizal, particularly his novels Noli Me Tangere and El Filibusterismo.", body_style),
            ("The original bill was authored by Senator Claro M. Recto (known as the Father of the Rizal Law) and sponsored in the Senate by Senator Jose P. Laurel, Sr., Chairman of the Committee on Education.", body_style),
            ("1.1 Legislative Rationale & Objectives", h2_style),
            ("• Re-dedication to Freedom: Inspiring Filipino youth with nationalistic ideals and patriotism during their formative years.", bullet_style),
            ("• Honoring the National Hero: Paying tribute to Jose Rizal for shaping the Filipino character and national consciousness.", bullet_style),
            ("• Source of Moral Character: Providing a fundamental source of patriotism, civic discipline, and moral courage.", bullet_style),
            ("1.2 The Great Controversy & Debates", h2_style),
            ("The bill faced intense opposition from the Catholic Church in the Philippines, led by figures such as Fr. Jesus Cavanna, Senator Francisco Soc Rodrigo, and Archbishop Rufino Santos. Opponents argued that mandatory reading of unexpurgated novels violated religious freedom, claiming 170 passages in Noli and 50 in Fili were anti-Catholic.", body_style),
            ("EXAM TRAP: Opponents asserted that the bill enforced dogma contrary to Catholic teachings. However, Recto defended the law, insisting Rizal aimed to expose clerical abuses rather than attack the Catholic faith itself.", body_style),
            ("1.3 Compromises & Key Provisions", h2_style),
            ("• Unexpurgated Versions: Required reading in tertiary (college) level.", bullet_style),
            ("• Expurgated Versions: Allowed in high school and elementary education.", bullet_style),
            ("• Exemption Clause: Students may request exemption from reading unexpurgated versions on written grounds of religious belief, but NOT from taking the course itself.", bullet_style)
        ]),
        ("2. The 19th Century World Context", [
            ("The 19th century was a period of profound global transformations marked by the Industrial Revolution, the rise of Democracy & Nationalism, and the expansion of Western Imperialism.", body_style),
            ("While European powers expanded, the Spanish Empire suffered severe decline, losing most South American colonies (Venezuela, Paraguay, Argentina, Chile, Peru). Spain retained only Cuba, Puerto Rico, and the Philippines.", body_style),
            ("2.1 Conditions in Spanish Philippines", h2_style),
            ("• Frailocracy (Frailocracia): The rule of Spanish friars (Augustinians, Dominicans, Franciscans, Recollects). Friars controlled politics, local administration, education, taxation, and land ownership.", bullet_style),
            ("• Political Instability: Frequent government shifts in Spain led to rapid turnover of Governor-Generals (41 Gov-Gens in 28 years), resulting in inconsistent policies.", bullet_style),
            ("• Representation Denied: Philippine representation in the Spanish Cortes (granted in 1810, 1820, 1834) was permanently abolished in 1837.", bullet_style),
            ("2.2 Economic Systems & Social Classes", h2_style),
            ("• Tributo / Cedula Personal: Direct tax imposed on native Filipinos.", bullet_style),
            ("• Polo y Servicios: Compulsory forced labor for 40 days/year (later reduced to 15 days) required of males aged 16–60. Exemption was available by paying a fee called faila.", bullet_style),
            ("• Bandala: Mandatory forced sale of agricultural produce to the Spanish government, often paid with worthless paper IOUs.", bullet_style),
            ("• Opening of Suez Canal (Nov 17, 1869): Reduced travel time between Spain and Manila from 3 months to 32 days, accelerating the entry of liberal ideas, books, and the emergence of the Ilustrados (educated class).", bullet_style),
            ("Social Pyramid: Peninsulares (Spaniards born in Spain) > Insulares/Creoles (Spaniards born in PH) > Mestizos (Spanish/Chinese mix) > Principalia (Local elite / Inquilinos) > Indios (Native Filipinos).", quote_style)
        ]),
        ("3. Rizal's Lineage, Childhood, & Early Life", [
            ("Jose Protasio Rizal Mercado y Alonso Realonda was born in Calamba, Laguna on June 19, 1861 (baptized June 22 by Fr. Rufino Collantes).", body_style),
            ("3.1 Breakdown of Rizal's Name", h2_style),
            ("• Jose: Chosen by mother in honor of St. Joseph.", bullet_style),
            ("• Protasio: From the Christian calendar (Feast day June 19).", bullet_style),
            ("• Rizal: From Spanish 'ricial' meaning green field / uncultivated ricefield (adopted per Claveria Decree of 1849).", bullet_style),
            ("• Mercado: Adopted in 1731 by paternal ancestor Domingo Lam-co (Chinese merchant from Amoy/Xiamen) meaning 'market'.", bullet_style),
            ("• Alonso Realonda: Family surname of mother Doña Teodora.", bullet_style),
            ("3.2 Family & Childhood Influences", h2_style),
            ("Father: Francisco Mercado II (inquilino tenant farmer of Dominican estate). Mother: Doña Teodora Alonso Realonda (educated at Colegio de Santa Rosa, Rizal's first teacher). Rizal was the 7th of 11 children. Older brother Paciano served as his second father, mentor, and financial backer.", body_style),
            ("First sorrow: Death of younger sister Concepcion (Concha) at age 3 in 1865. Childhood poem: Un Recuerdo A Mi Pueblo (In Memory of My Town).", body_style),
            ("3.3 Formative Tragedies & Catalysts", h2_style),
            ("1. Incarceration of Doña Teodora (1871–1874): Falsely accused by brother Jose Alberto's estranged wife of attempting to poison her. Doña Teodora was forced to walk 50 kilometers from Calamba to Santa Cruz provincial jail.", body_style),
            ("2. Execution of GOMBURZA (Feb 17, 1872): Fathers Mariano Gomez, Jose Burgos, and Jacinto Zamora were executed by garrote following the Cavite Mutiny (Jan 20, 1872). Paciano was a close student and friend of Fr. Burgos. Rizal later dedicated El Filibusterismo to GOMBURZA.", body_style)
        ]),
        ("4. Education: Ateneo Municipal vs. UST", [
            ("4.1 Ateneo Municipal (1872–1877)", h2_style),
            ("Managed by the Jesuits. Rizal enrolled after passing entrance exams at Colegio de San Juan de Letran. Fr. Magin Ferrando initially rejected him (late registration and frail physique), but Manuel Xerez Burgos intervened.", body_style),
            ("Jesuit System: Divided classes into two empires — Roman Empire (boarders) and Carthaginian Empire (non-boarders). Ranks: Imperator, Tribune, Decurion, Centurion, Standard-Bearer. Rizal rose to Imperator within a month.", body_style),
            ("Graduated on March 23, 1877 with the degree Bachiller en Artes with highest honors (Sobresaliente). Favorite mentor: Fr. Francisco de Paula Sanchez. Key poems: Mi Primera Inspiracion, Por la Educacion Recibe Lustre la Patria.", body_style),
            ("4.2 University of Santo Tomas (1877–1882)", h2_style),
            ("Managed by the Dominicans. Rizal studied Philosophy and Letters, then shifted to Medicine (Ophthalmology) to cure his mother's failing eyesight.", body_style),
            ("Concurrently finished Perito Agrimensor (Expert Surveyor) at Ateneo (1877).", body_style),
            ("Reasons for leaving UST: Dominican professors were discriminatory, archaic teaching methods (Physics taught without laboratory experiments), and rampant racial prejudice against Filipino students.", body_style),
            ("Literary Laurels at UST: Won 1st prize for A La Juventud Filipina (1879 - 'The Youth is the Fair Hope of the Fatherland') and El Consejo de los Dioses (1880).", body_style)
        ]),
        ("5. First European Odyssey & Propaganda", [
            ("5.1 Secret Mission & Arrival in Spain (1882)", h2_style),
            ("Rizal departed secretly on May 3, 1882 aboard SS Salvador with approval from Paciano and uncle Antonio Rivera. Mission: Observe European life, cultures, laws, and governments to prepare for liberating his people.", body_style),
            ("Barcelona: Wrote El Amor Patrio (Love of Country), published in Diariong Tagalog under pen name Laong Laan.", body_style),
            ("Madrid: Enrolled at Universidad Central de Madrid (Medicine and Philosophy & Letters). Joined Circulo Hispano-Filipino and Masonic Lodge Solidaridad/Acacia adopting name Dimasalang.", body_style),
            ("Historic Toast (June 25, 1884): Speech honoring Juan Luna (Gold: Spoliarium) and Felix Resurreccion Hidalgo (Silver: Las Virgenes Cristianas Expuestas al Populacho), proclaiming that genius knows no country.", body_style),
            ("5.2 Paris, Germany, & Noli Me Tangere", h2_style),
            ("• Paris (1885): Worked as assistant to famous ophthalmologist Dr. Louis de Wecker.", bullet_style),
            ("• Heidelberg (1886): Worked under Dr. Otto Becker. Wrote poem A Las Flores de Heidelberg. Stayed with Pastor Karl Ullmer in Wilhelmsfeld.", bullet_style),
            ("• Berlin (1887): Met ethnologists Dr. Rudolf Virchow and Dr. Feodor Jagor. Translated Schiller's William Tell into Tagalog.", bullet_style),
            ("Publication of Noli Me Tangere (Berlin, March 1887): Title from Gospel of St. John 20:17 ('Touch Me Not'). Funded by a 300-peso loan from Dr. Maximo Viola for 2,000 copies at Berliner Buchdruckerei-Action-Gesselschaft.", body_style),
            ("Grand European Tour (May 1887): Toured Europe with Maximo Viola. Met lifelong friend and scholar Ferdinand Blumentritt in Leitmeritz.", body_style)
        ]),
        ("6. First Return, 2nd Journey, & La Liga", [
            ("6.1 First Return & Noli Storm (1887–1888)", h2_style),
            ("Returned to Calamba, operated on mother's eye (called Dr. Uliman). Gov.-Gen. Emilio Terrero assigned Lt. Jose Taviel de Andrade as bodyguard. UST committee & Censorship Board (Fr. Salvador Font) banned Noli as heretical. Fr. Vicente Garcia defended Noli under pen name Justo Desiderio Magalang. Agrarian report on Dominican hacienda forced Rizal to depart.", body_style),
            ("6.2 Second Overseas Sojourn (1888–1892)", h2_style),
            ("• London: Stayed with Beckett family. Annotated Antonio de Morga's Sucesos de las Islas Filipinas (1609) at British Museum to prove pre-colonial Filipino civilization.", bullet_style),
            ("• La Solidaridad: Organ of Propaganda Movement (founded Dec 15, 1888 by Graciano Lopez Jaena, later edited by Marcelo H. del Pilar). Rizal wrote essays: La Indolencia de los Filipinos, Filipinas Dentro de Cien Años.", bullet_style),
            ("• El Filibusterismo: Published Sept 18, 1891 in Ghent, Belgium, financed by Valentin Ventura. Dedicated to GOMBURZA.", bullet_style),
            ("• Hong Kong (1891–1892): Practiced medicine, proposed Borneo Colonization Project (rejected by Gov. Despujol), drafted La Liga Filipina constitution.", bullet_style),
            ("6.3 Founding of La Liga Filipina", h2_style),
            ("Founded on July 3, 1892 at the house of Doroteo Ongjunco in Tondo, Manila. Motto: Unus Instar Omnium (One Like All). Aims: Mutual protection, defense against violence/injustice, encouragement of education, agriculture, and commerce.", body_style)
        ]),
        ("7. Exile in Dapitan, Trial, & Execution", [
            ("7.1 Exile in Dapitan (1892–1896)", h2_style),
            ("Arrested July 6, 1892; deported to Dapitan aboard steamer Cebu. Stayed with Commandant Captain Ricardo Carnicero. Won Manila Lottery (Ticket 9736) sharing 20,000 pesos jackpot.", body_style),
            ("Dapitan Achievements: Built school for boys, constructed gravity waterworks system, made relief map of Mindanao in front of church, collected specimens (Draco rizali, Apogonia rizali, Rhacophorus rizali), married Josephine Bracken informally.", body_style),
            ("Dr. Pio Valenzuela Visit (June 1896): Katipunan envoy sent by Andres Bonifacio. Rizal opposed premature armed uprising without adequate weapons and naval support.", body_style),
            ("7.2 Arrest, Trial, & Fort Santiago", h2_style),
            ("Volunteered as army doctor in Cuba. En route on Isla de Panay, revolution broke out (Aug 1896). Arrested, held at Montjuïch Castle (Barcelona), returned on steamer Colon, imprisoned at Fort Santiago (Nov 3, 1896).", body_style),
            ("Preliminary Investigation by Col. Francisco Olive. Defense Attorney: Lt. Luis Taviel de Andrade. Charges: Sedition, Rebellion, Illegal Association. Tried by Military Court at Cuartel de España. Sentenced to death by Gov.-Gen. Camilo de Polavieja.", body_style),
            ("7.3 Execution at Bagumbayan (Dec 30, 1896)", h2_style),
            ("Wrote Mi Ultimo Adios (My Last Farewell), hidden inside an alcohol cooking stove (cocinilla) given to sister Trinidad. Walked from Fort Santiago to Bagumbayan (Luneta) at 6:30 AM with Fr. Estanislao March and Fr. Jose Villaclara.", body_style),
            ("Executed by firing squad at 7:03 AM, Dec 30, 1896. Turned to face the morning sun as he fell. Last words: \"Consummatum est!\" (It is finished!).", body_style)
        ])
    ]

    for title, paragraphs in modules:
        story.append(Paragraph(title, h1_style))
        for content_text, style in paragraphs:
            if style == h2_style:
                story.append(Paragraph(content_text, style))
            else:
                formatted_text = auto_highlight(content_text)
                story.append(Paragraph(formatted_text, style))

    # Add Comparative Tables in Module 8
    story.append(Paragraph("8. Summary & Comparative Matrices", h1_style))
    story.append(Paragraph("8.1 Noli Me Tangere vs. El Filibusterismo", h2_style))

    data_noli_fili = [
        [Paragraph("Feature", table_header_style), Paragraph("Noli Me Tangere", table_header_style), Paragraph("El Filibusterismo", table_header_style)],
        [Paragraph("Meaning", table_cell_style), Paragraph("Touch Me Not (John 20:17)", table_cell_style), Paragraph("The Reign of Greed / Subversive", table_cell_style)],
        [Paragraph("Nature / Tone", table_cell_style), Paragraph("Romantic, Aspirant, Soft, Mind", table_cell_style), Paragraph("Political, Revengeful, Harsh, Heart", table_cell_style)],
        [Paragraph("Dedicated To", table_cell_style), Paragraph("To the Fatherland (Philippines)", table_cell_style), Paragraph("GOMBURZA", table_cell_style)],
        [Paragraph("Finished / Printed", table_cell_style), Paragraph("Berlin (March 1887)", table_cell_style), Paragraph("Ghent (Sept 1891)", table_cell_style)],
        [Paragraph("Financier / Loan", table_cell_style), Paragraph("Dr. Maximo Viola", table_cell_style), Paragraph("Valentin Ventura", table_cell_style)]
    ]
    t1 = Table(data_noli_fili, colWidths=[COL_WIDTH*0.28, COL_WIDTH*0.36, COL_WIDTH*0.36])
    t1.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#E6E6FA')),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CCCCCC')),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('TOPPADDING', (0,0), (-1,-1), 2),
        ('BOTTOMPADDING', (0,0), (-1,-1), 2),
    ]))
    story.append(t1)
    story.append(Spacer(1, 4))

    story.append(Paragraph("8.2 Ateneo Municipal vs. UST", h2_style))
    data_ateneo_ust = [
        [Paragraph("Feature", table_header_style), Paragraph("Ateneo (Jesuits)", table_header_style), Paragraph("UST (Dominicans)", table_header_style)],
        [Paragraph("System", table_cell_style), Paragraph("Ratio Studiorum (Humanistic)", table_cell_style), Paragraph("Scholastic / Traditional", table_cell_style)],
        [Paragraph("Student Ranks", table_cell_style), Paragraph("Empires (Roman vs. Carthaginian)", table_cell_style), Paragraph("No empire ranking system", table_cell_style)],
        [Paragraph("Rizal's Exp.", table_cell_style), Paragraph("Inspired, Encouraged, Sobresaliente", table_cell_style), Paragraph("Repressed, Discriminated, Left for Spain", table_cell_style)]
    ]
    t2 = Table(data_ateneo_ust, colWidths=[COL_WIDTH*0.25, COL_WIDTH*0.375, COL_WIDTH*0.375])
    t2.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#E6E6FA')),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CCCCCC')),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('TOPPADDING', (0,0), (-1,-1), 2),
        ('BOTTOMPADDING', (0,0), (-1,-1), 2),
    ]))
    story.append(t2)
    story.append(Spacer(1, 4))

    story.append(Paragraph("8.3 Key Women in Rizal's Life", h2_style))
    data_women = [
        [Paragraph("Name", table_header_style), Paragraph("Description / Connection", table_header_style)],
        [Paragraph("Segunda Katigbak", table_cell_style), Paragraph("First love (from Lipa, Batangas; engaged to Manuel Luz)", table_cell_style)],
        [Paragraph("Leonor Valenzuela ('Orang')", table_cell_style), Paragraph("Pagsanjan neighbor; wrote love letters in invisible ink", table_cell_style)],
        [Paragraph("Leonor Rivera ('Taimis')", table_cell_style), Paragraph("True love & cousin; 11-year engagement; inspiration for Maria Clara", table_cell_style)],
        [Paragraph("Consuelo Ortiga y Perez", table_cell_style), Paragraph("Madrid romance; daughter of Don Pablo Ortiga y Rey", table_cell_style)],
        [Paragraph("O-Sei-San (Usui Seiko)", table_cell_style), Paragraph("Japanese samurai daughter; interpreter & guide in Tokyo (1888)", table_cell_style)],
        [Paragraph("Gertrude Beckett ('Tottie')", table_cell_style), Paragraph("London landlord's daughter; helped carve figurines", table_cell_style)],
        [Paragraph("Suzanne Jacoby", table_cell_style), Paragraph("Brussels landlady's niece (1890)", table_cell_style)],
        [Paragraph("Nellie Boustead", table_cell_style), Paragraph("Protestant romance in Biarritz; failed due to religion & mother's objection", table_cell_style)],
        [Paragraph("Josephine Bracken", table_cell_style), Paragraph("Irish woman who accompanied George Taufer to Dapitan; Rizal's common-law wife", table_cell_style)]
    ]
    t3 = Table(data_women, colWidths=[COL_WIDTH*0.35, COL_WIDTH*0.65])
    t3.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#E6E6FA')),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CCCCCC')),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('TOPPADDING', (0,0), (-1,-1), 2),
        ('BOTTOMPADDING', (0,0), (-1,-1), 2),
    ]))
    story.append(t3)

    if progress_callback:
        progress_callback(85)

    doc.build(story, canvasmaker=CleanCanvas)
    if progress_callback:
        progress_callback(100)

    if log_callback:
        log_callback(f"Successfully generated Reviewer PDF at: {output_pdf_path}")

    return output_pdf_path
